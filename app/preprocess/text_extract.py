import fitz
import pptx
from docx import Document

def extract_text(file_path, file_name):
    """Extracts text from different file types."""
    ext = file_name.split('.')[-1].lower()
    
    try:
        if ext == "pdf":
            return extract_text_from_pdf(file_path)
        elif ext in ["doc", "docx"]:
            return extract_text_from_docx(file_path)
        elif ext in ["ppt", "pptx"]:
            return extract_text_from_pptx(file_path)
        elif ext == "txt":
            return extract_text_from_txt(file_path)
        else:
            return "Unsupported file format"
    except Exception as e:
        print(f"Error extracting text from {file_name}: {e}")
        return "Error extracting text"

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file."""
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text("text") + "\n"
    return text.strip()

def extract_text_from_docx(file_path):
    """Extracts text from a DOCX file."""
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    """Extracts text from a PowerPoint file."""
    ppt = pptx.Presentation(file_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_txt(file_path):
    """Extracts text from a TXT file."""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read().strip()
